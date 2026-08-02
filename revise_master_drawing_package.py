from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Pt

SOURCE = Path(r"F:\OMNI\Vern\AQUILA_Sovereign_Master_Drawing_Package_REV_B- drat.pptx")
OUTPUT = Path(r"F:\OMNI\Vern\AQUILA_Sovereign_Master_Drawing_Package_REV_B_CORRECTED.pptx")

NAVY = RGBColor(18, 55, 105)
BLUE = RGBColor(38, 102, 170)
LIGHT_BLUE = RGBColor(231, 240, 248)
GREEN = RGBColor(31, 122, 81)
LIGHT_GREEN = RGBColor(231, 246, 238)
ORANGE = RGBColor(222, 116, 38)
LIGHT_ORANGE = RGBColor(252, 239, 226)
PURPLE = RGBColor(104, 76, 150)
LIGHT_PURPLE = RGBColor(241, 235, 248)
RED = RGBColor(184, 48, 48)
LIGHT_RED = RGBColor(252, 235, 235)
GRAY = RGBColor(100, 106, 112)
LIGHT_GRAY = RGBColor(244, 245, 246)
DARK = RGBColor(31, 35, 39)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(190, 196, 201)


def pt(value):
    return Pt(value)


def set_text(shape, text, size=12, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font="Arial", valign=MSO_ANCHOR.MIDDLE, margin=4):
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True
    shape.text_frame.vertical_anchor = valign
    shape.text_frame.margin_left = pt(margin)
    shape.text_frame.margin_right = pt(margin)
    shape.text_frame.margin_top = pt(2)
    shape.text_frame.margin_bottom = pt(2)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def text_box(slide, x, y, w, h, text, size=12, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font="Arial", fill=None, line=None, radius=False,
             valign=MSO_ANCHOR.MIDDLE, margin=4):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, pt(x), pt(y), pt(w), pt(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = pt(0.8)
    return set_text(shape, text, size, bold, color, align, font, valign, margin)


def line(slide, x1, y1, x2, y2, color=BORDER, width=1.0, arrow=False):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, pt(x1), pt(y1), pt(x2), pt(y2))
    connector.line.color.rgb = color
    connector.line.width = pt(width)
    if arrow:
        connector.line.end_arrowhead = True
    return connector


def remove_all_shapes(slide):
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def replace_native_text(slide, old_text, new_text):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip() == old_text:
            shape.text_frame.text = new_text
            return True
    return False


def add_sheet_header(slide, number, drawing_no, title):
    text_box(slide, 42, 52, 79, 28, "", fill=RGBColor(238, 244, 248), line=None)
    text_box(slide, 49, 56, 53, 21, "SHEET", 11, False, RGBColor(55, 82, 105))
    text_box(slide, 96, 54, 28, 26, f"{number:02d}", 13, False, RGBColor(55, 82, 105))
    text_box(slide, 139, 42, 145, 43, drawing_no, 18, True, DARK)
    text_box(slide, 290, 50, 450, 31, title, 14, False, RGBColor(85, 85, 85))
    text_box(slide, 1052, 57, 381, 22, "AQUILA SOVEREIGN - MASTER DRAWING PACKAGE REV B", 8,
             False, RGBColor(130, 130, 130), PP_ALIGN.RIGHT)


def add_legacy_override(slide, sheet_no):
    text_box(slide, 1208, 691, 180, 24,
             f"CONTROLLED OVERRIDE  |  SHEET {sheet_no:02d} OF 10  |  REV B",
             7.2, True, NAVY, PP_ALIGN.CENTER, fill=WHITE, line=NAVY, radius=False, margin=2)


def add_patent_ribbon(slide, patent_text):
    text_box(slide, 690, 78, 698, 17, patent_text, 7.2, True, GREEN,
             PP_ALIGN.RIGHT, fill=LIGHT_GREEN, line=GREEN, radius=True, margin=5)


def add_tree_node(slide, x, y, w, h, title, subtitle="", color=NAVY, fill=LIGHT_BLUE, size=9):
    text = title if not subtitle else f"{title}\n{subtitle}"
    return text_box(slide, x, y, w, h, text, size, True, color, PP_ALIGN.CENTER,
                    fill=fill, line=color, radius=True, margin=4)


def create_sheet_08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(245, 246, 247)
    add_sheet_header(slide, 8, "DWG-MECH-08", "Power + Data Trees")
    text_box(slide, 42, 96, 1356, 672, "", fill=WHITE, line=BORDER)

    text_box(slide, 60, 108, 650, 28, "POWER TREE - 48 V HYBRID BUS", 14, True, NAVY,
             PP_ALIGN.LEFT, fill=LIGHT_BLUE, line=NAVY)
    text_box(slide, 730, 108, 650, 28, "DATA TREE - JETSON AGX ORIN HUB", 14, True, NAVY,
             PP_ALIGN.LEFT, fill=LIGHT_BLUE, line=NAVY)

    add_tree_node(slide, 270, 145, 230, 42, "48 V FIREFLY HYBRID BUS", "AWG12 | 30 A fuse | Molex MX150L")
    power_nodes = [
        (70, 220, 145, 54, "DIRECT 48 V", "FOC drill driver\n150 W peak | XT60", RED, LIGHT_RED),
        (230, 220, 170, 54, "VICOR DC-DC", "48 V -> 12 V + 5 V\nDFMA REV B", BLUE, LIGHT_BLUE),
        (415, 220, 145, 54, "24 V DIRECT", "Venturi + pump\nAWG16 | 13 W", GREEN, LIGHT_GREEN),
        (575, 220, 115, 54, "5 V LOGIC", "Distributed\nsensors | AWG22", PURPLE, LIGHT_PURPLE),
    ]
    for x, y, w, h, title, sub, color, fill in power_nodes:
        add_tree_node(slide, x, y, w, h, title, sub, color, fill, 8.2)
        line(slide, 385, 187, x + w / 2, y, color, 1.4, True)

    add_tree_node(slide, 230, 304, 78, 40, "12 V", "8 A", BLUE, LIGHT_BLUE, 8)
    add_tree_node(slide, 322, 304, 78, 40, "5 V", "4 A", PURPLE, LIGHT_PURPLE, 8)
    line(slide, 315, 274, 269, 304, BLUE, 1.2, True)
    line(slide, 315, 274, 361, 304, PURPLE, 1.2, True)

    rail12 = ["Jetson 65 W", "AERIS-10 25 W", "FLIR 5 W", "EM clutch 15 W", "Heater 10 W", "TMC2209 8 W"]
    rail5 = ["STM32H7 2 W", "SGP30/SHT31 0.5 W", "ReSpeaker 1 W", "LIDAR 0.5 W"]
    text_box(slide, 70, 365, 315, 132, "12 V LOADS\n" + "\n".join(f"- {item}" for item in rail12),
             8.3, False, DARK, fill=LIGHT_GRAY, line=BORDER, valign=MSO_ANCHOR.TOP, margin=8)
    text_box(slide, 400, 365, 290, 132, "5 V / AUXILIARY LOADS\n" + "\n".join(f"- {item}" for item in rail5),
             8.3, False, DARK, fill=LIGHT_GRAY, line=BORDER, valign=MSO_ANCHOR.TOP, margin=8)

    text_box(slide, 70, 515, 620, 75,
             "POWER BUDGET\nContinuous design load: 212 W   |   Peak design load: 332 W\nAvailable budget: 480 W   |   Peak margin: +148 W   |   Continuous margin: +268 W",
             9, True, GREEN, fill=LIGHT_GREEN, line=GREEN, radius=True, margin=8)
    text_box(slide, 70, 600, 620, 60,
             "CONTROL NOTES\n1. Branch protection and connector ratings require verification before release.\n2. Shield returns terminate at controlled chassis ground; avoid ground loops.",
             7.8, False, DARK, fill=WHITE, line=BORDER, valign=MSO_ANCHOR.TOP, margin=8)

    add_tree_node(slide, 940, 145, 230, 42, "JETSON AGX ORIN 64 GB", "Central compute / ROS 2 / AI", NAVY, LIGHT_BLUE)
    data_nodes = [
        (750, 225, 125, 55, "USB-C / USB3", "MinION | ZED 2i\nAERIS-10 FT601", BLUE, LIGHT_BLUE),
        (890, 225, 115, 55, "MIPI CSI-2", "Sony ILX-LR1\nFLIR Hadron", PURPLE, LIGHT_PURPLE),
        (1020, 225, 115, 55, "CAN", "STM32H7 hub\nU2D2 / servos", GREEN, LIGHT_GREEN),
        (1150, 225, 105, 55, "SERIAL", "UART / I2C\nSPI", ORANGE, LIGHT_ORANGE),
        (1270, 225, 100, 55, "MAVLink", "Autopilot\nADS-B", GRAY, LIGHT_GRAY),
    ]
    for x, y, w, h, title, sub, color, fill in data_nodes:
        add_tree_node(slide, x, y, w, h, title, sub, color, fill, 7.8)
        line(slide, 1055, 187, x + w / 2, y, color, 1.4, True)

    data_details = [
        (750, 310, 195, 112, "HIGH-BANDWIDTH\n- MinION Mk1C: USB-C\n- ZED 2i: USB-C\n- AERIS-10: USB3 / FT601\n- Sony: MIPI 4-lane\n- FLIR: MIPI 2-lane", BLUE, LIGHT_BLUE),
        (960, 310, 195, 112, "CONTROL BUSES\n- STM32H7: CAN 1 Mbit/s\n- RFID / SHT31 / valves\n- U2D2 -> J1 -> J4\n- Drill telemetry: UART\n- Sensors: I2C / SPI", GREEN, LIGHT_GREEN),
        (1170, 310, 200, 112, "FLIGHT INTERFACE\n- MAVLink UART\n- Autopilot state\n- ADS-B transponder\n- Time synchronization\n- Health / fault status", ORANGE, LIGHT_ORANGE),
    ]
    for x, y, w, h, text, color, fill in data_details:
        text_box(slide, x, y, w, h, text, 7.8, False, DARK, fill=fill, line=color,
                 valign=MSO_ANCHOR.TOP, margin=8)

    text_box(slide, 750, 445, 620, 65,
             "PROTOCOL LEGEND\nUSB / FT601: blue   |   CAN: green   |   UART / MAVLink: orange\nMIPI: purple   |   I2C / SPI: gray   |   All external links require labeled service loops.",
             8, True, DARK, fill=LIGHT_GRAY, line=BORDER, radius=True, margin=8)
    text_box(slide, 750, 525, 620, 91,
             "EMI / EMC MITIGATION - PAT-20 ELECTRONIC MODULE INTEGRATION\n- Shielded FPC and copper-foil ground for analog signal paths.\n- Ferrite cores placed at measured cable resonance antinodes.\n- Segmented ground plane with polarization channels; chassis bond at controlled points.\n- Frequency-hopping PWM retained as a validation item, not a released compliance claim.",
             8, False, DARK, fill=LIGHT_ORANGE, line=ORANGE, radius=True,
             valign=MSO_ANCHOR.TOP, margin=8)

    # Editable title block.
    text_box(slide, 750, 635, 620, 102, "", fill=WHITE, line=DARK)
    line(slide, 750, 665, 1370, 665, DARK, 0.8)
    line(slide, 750, 700, 1370, 700, DARK, 0.8)
    line(slide, 930, 635, 930, 737, DARK, 0.8)
    line(slide, 1190, 635, 1190, 737, DARK, 0.8)
    text_box(slide, 758, 640, 165, 20, "AQUILA GEOLOGICAL SYSTEMS", 8, True, NAVY)
    text_box(slide, 938, 640, 245, 20, "DWG-MECH-08", 10, True, DARK)
    text_box(slide, 1198, 640, 164, 20, "SHEET 08 OF 10 | REV B", 8, True, DARK, PP_ALIGN.CENTER)
    text_box(slide, 758, 670, 165, 25, "SCALE: NTS\nUNITS: SI / mm", 7, False, DARK)
    text_box(slide, 938, 670, 245, 25, "POWER + DATA TREES\nCONTROLLED ENGINEERING SCHEMATIC", 8, True, DARK)
    text_box(slide, 1198, 670, 164, 25, "DATE: 2026-07\nSTATUS: VALIDATION", 7, False, DARK, PP_ALIGN.CENTER)
    text_box(slide, 758, 705, 604, 27,
             "NOT FOR MANUFACTURING RELEASE - VERIFY CONNECTOR PINOUTS, FUSING, GROUNDING, AND LOAD CASES",
             7, True, RED, PP_ALIGN.CENTER)
    return slide


def rebuild_directory(slide):
    remove_all_shapes(slide)
    text_box(slide, 68, 68, 1435, 47, "Sheet Directory - Corrected 10-Sheet Set", 24, False, DARK)
    cols = [(68, 82, "SHEET"), (150, 195, "DRAWING NO."), (345, 502, "TITLE"), (847, 526, "CONTENT")]
    y = 140
    for x, w, title in cols:
        text_box(slide, x, y, w, 20, title, 9, True, GRAY)
        line(slide, x, y + 19, x + w, y + 19, NAVY, 1.2)
    rows = [
        ("01", "DWG-MECH-01", "Tier 1 Airframe Assembly", "2200 mm span, hub, motor mounts, rail clamps, GD&T"),
        ("02", "DWG-MECH-02", "OmniSpectral CF Housing STR-001", "Ø300 x 270 mm shell, V-groove rail, EMI gaskets"),
        ("03", "DWG-MECH-03", "Sensor Pod Internal Layout", "15 instruments, zones, clearances, service envelopes"),
        ("04", "DWG-MECH-04", "DrACO Chassis + Drill + Carousel", "PAT-01/02/03/10/11/12/17/18; DFMA REV B"),
        ("05", "DWG-MECH-05", "Pneumatic Transport + Bio Enclosure", "PAT-04/05/06/13/14/15/16; validation windows"),
        ("06", "DWG-MECH-06", "Robotic Arm + End-Effector", "PAT-07/08/09/19; 4-DOF, 280 mm reach"),
        ("07", "DWG-MECH-07", "Master Interface + BOM Summary", "PAT-20; power, data, pneumatic, thermal paths"),
        ("08", "DWG-MECH-08", "Power + Data Trees", "48 V distribution, load budget, protocols, EMI controls"),
        ("09", "DWG-MECH-09", "Pneumatic P&ID + DrACO Sequence", "ISO 1219 schematic, six-step flow, RLVR matrix"),
        ("10", "DWG-MECH-10", "Master BOM + Compliance + Roadmap", "BOM, NASARAP checks, patent index, deployment phases"),
    ]
    row_y = 165
    row_h = 49
    for idx, row in enumerate(rows):
        fill = WHITE if idx % 2 == 0 else RGBColor(248, 249, 250)
        text_box(slide, 68, row_y, 1305, row_h, "", fill=fill, line=BORDER)
        text_box(slide, 76, row_y + 8, 58, 31, row[0], 12, True, NAVY, PP_ALIGN.CENTER,
                 fill=LIGHT_BLUE, line=None, radius=True)
        text_box(slide, 155, row_y + 5, 185, 38, row[1], 11, True, DARK)
        text_box(slide, 350, row_y + 5, 490, 38, row[2], 11, False, DARK)
        text_box(slide, 852, row_y + 5, 515, 38, row[3], 9.5, False, GRAY)
        row_y += row_h
    text_box(slide, 68, 675, 1305, 56,
             "REV B PACKAGE CORRECTION 01: Added missing DWG-MECH-08; corrected set count; normalized sheet numbering; "
             "added patent traceability; retained source artwork as controlled raster references.",
             10, True, GREEN, fill=LIGHT_GREEN, line=GREEN, radius=True, margin=10)


def main():
    prs = Presentation(SOURCE)

    # Cover corrections.
    replace_native_text(prs.slides[0], "AQUILA Sovereign - Mechanical Drawing Set, Rev B",
                        "AQUILA Sovereign - Mechanical Drawing Set, Rev B - Package Correction 01")
    replace_native_text(prs.slides[0], "9 of Set", "10 of Set")
    replace_native_text(prs.slides[0], "REV B", "REV B / CORR 01")

    rebuild_directory(prs.slides[1])

    # Legacy title-block count correction for existing Sheets 1-7.
    for sheet_no, slide_index in enumerate(range(2, 9), 1):
        add_legacy_override(prs.slides[slide_index], sheet_no)

    # Patent traceability ribbons.
    add_patent_ribbon(prs.slides[5], "PATENTS: PAT-01 / 02 / 03 / 10 / 11 / 12 / 17 / 18")
    add_patent_ribbon(prs.slides[6], "PATENTS: PAT-04 / 05 / 06 / 13 / 14 / 15 / 16")
    add_patent_ribbon(prs.slides[7], "PATENTS: PAT-07 / 08 / 09 / 19")
    add_patent_ribbon(prs.slides[8], "PATENT: PAT-20 ELECTRONIC MODULE INTEGRATION")

    # Correct Sheet 10 artwork alignment.
    for shape in prs.slides[10].shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape.left = pt(42.4)
            shape.top = pt(101.5)
            break

    # Add Sheet 08 and insert it between existing Sheets 07 and 09.
    new_slide = create_sheet_08(prs)
    slide_ids = prs.slides._sldIdLst
    new_id = slide_ids[-1]
    slide_ids.remove(new_id)
    slide_ids.insert(9, new_id)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)} (cover + directory + 10 engineering sheets)")


if __name__ == "__main__":
    main()
